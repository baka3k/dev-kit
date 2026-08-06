#!/usr/bin/env python3
"""Create a fictional, explicitly illustrative deck on the hi-pptx engine.

The deck exercises the high-contrast light minimalist system end to end:
Palette A "Crisp Swiss", the four strict layouts (hero metric, stark
editorial split, max-3 tiles, minimal quote), 60-30-10 light color
discipline, titles under 35 characters, and body blocks under 40 words.
It must pass lint_pptx.py with zero errors and zero warnings.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


# Palette A "Crisp Swiss" tokens. Badge tint is the accent at 8% over white.
C = {
    "bg": "FFFFFF",
    "card": "F1F5F9",
    "primary": "020617",
    "secondary": "64748B",
    "accent": "4F46E5",
    "hairline": "E2E8F0",
    "obsidian": "0F172A",
    "tint": "F1F0FD",
    "white": "FFFFFF",
}

FONT = "Inter"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(C[color])


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color="primary",
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(C[color])
    return box


def add_box(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill="card",
    line="hairline",
    no_line=False,
    radius=True,
    corner_ratio=0.05,
    line_width=1,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # python-pptx writes a p:style whose effectRef points at the theme's
    # effect styles; the bundled default theme has no effectStyleList, and
    # LibreOffice then falls back to a drop shadow for ANY effectRef index.
    # The Swiss light engine keeps cards flat (tint fills, hairlines, accent
    # borders — never shadows), so drop the style element entirely. Fill,
    # line, and run fonts are all set explicitly and do not need it.
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)
    if radius:
        shape.adjustments[0] = corner_ratio
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C[fill])
    if no_line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(C[line])
        shape.line.width = Pt(line_width)
    return shape


def add_badge_pill(slide, text, x, y, w=3.4, h=0.36):
    add_box(slide, x, y, w, h, fill="tint", no_line=True, radius=True, corner_ratio=0.5)
    add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=10.5,
        color="accent",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, number, number_color="secondary"):
    add_text(
        slide,
        "ILLUSTRATIVE PIPELINE TEST · FICTIONAL SCENARIO",
        0.833,
        7.06,
        5.0,
        0.2,
        size=9,
        color="secondary",
        bold=True,
    )
    add_text(
        slide,
        str(number),
        12.10,
        7.06,
        0.4,
        0.2,
        size=9,
        color=number_color,
        align=PP_ALIGN.RIGHT,
    )


def make_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "A Bounded AI Pilot for Support Ops"
    prs.core_properties.subject = "Smoke test for hi-pptx"
    prs.core_properties.author = "Hi-PPTX Engine"
    blank = prs.slide_layouts[6]

    # 1 — Quiet cover
    slide = prs.slides.add_slide(blank)
    set_background(slide, "bg")
    add_badge_pill(slide, "ILLUSTRATIVE · FICTIONAL SCENARIO", 0.833, 0.625, 3.6)
    add_text(
        slide,
        "A bounded AI pilot\nfor support ops",
        0.833,
        2.15,
        9.8,
        1.9,
        size=44,
        color="primary",
        bold=True,
    )
    add_box(slide, 0.833, 4.42, 1.2, 0.035, fill="accent", line="accent", radius=False)
    add_text(
        slide,
        "Eight weeks, one workflow, human approval on every send.",
        0.833,
        4.78,
        8.4,
        0.4,
        size=15,
        color="secondary",
    )
    add_text(
        slide,
        "Prepared for pipeline testing · August 2026",
        0.833,
        6.55,
        5.8,
        0.3,
        size=12,
        color="secondary",
    )

    # 2 — Layout 1: Hero Metric on vast white
    slide = prs.slides.add_slide(blank)
    set_background(slide, "bg")
    add_badge_pill(slide, "BASELINE · WEEKLY VOLUME", 0.833, 0.625, 3.0)
    add_text(
        slide,
        "43%",
        0.75,
        1.75,
        6.5,
        2.1,
        size=120,
        color="accent",
        bold=True,
    )
    add_text(
        slide,
        "of support requests repeat questions the team already answered.\nOne retrieval step can catch most of them.",
        0.833,
        4.35,
        8.6,
        0.9,
        size=15,
        color="secondary",
    )
    add_text(
        slide,
        "Illustrative numbers for pipeline testing",
        0.833,
        5.60,
        6.0,
        0.25,
        size=9,
        color="secondary",
    )
    add_footer(slide, 2)

    # 3 — Layout 2: Stark Editorial Split (50/50, obsidian block)
    slide = prs.slides.add_slide(blank)
    set_background(slide, "bg")
    add_badge_pill(slide, "CHAPTER 2 · CONTROL", 0.833, 0.85, 2.7)
    add_text(
        slide,
        "Humans keep\nthe final call",
        0.833,
        2.45,
        5.4,
        1.7,
        size=40,
        color="primary",
        bold=True,
    )
    add_text(
        slide,
        "The pilot drafts responses.\nExperts approve every send.",
        0.833,
        4.55,
        4.6,
        0.8,
        size=14,
        color="secondary",
    )
    add_box(slide, 6.667, 0, 6.666, 7.5, fill="obsidian", line="obsidian", radius=False)
    add_text(
        slide,
        "02",
        7.50,
        2.30,
        4.0,
        1.9,
        size=110,
        color="white",
        bold=True,
    )
    add_text(
        slide,
        "CONTROL FIRST",
        7.60,
        4.45,
        4.0,
        0.3,
        size=12,
        color="white",
        bold=True,
    )
    add_footer(slide, 3, number_color="white")

    # 4 — Layout 3: Three tiles, one featured
    slide = prs.slides.add_slide(blank)
    set_background(slide, "bg")
    add_text(
        slide,
        "Three guardrails, one pilot",
        0.833,
        0.72,
        8.0,
        0.7,
        size=33,
        color="primary",
        bold=True,
    )
    tiles = [
        ("01", "Scoped data", "Approved sources only; no raw customer records.", False),
        ("02", "Human approval", "Every draft is reviewed before it is sent.", True),
        ("03", "Clean exit", "Kill switch, audit trail, and a handback plan.", False),
    ]
    for i, (number, label, body, featured) in enumerate(tiles):
        x = 0.833 + i * 4.01
        add_box(
            slide,
            x,
            2.10,
            3.73,
            3.50,
            fill="white" if featured else "card",
            line="accent",
            no_line=not featured,
            radius=True,
            corner_ratio=0.045,
            line_width=1.5,
        )
        add_text(
            slide,
            number,
            x + 0.30,
            2.50,
            1.0,
            0.4,
            size=18,
            color="accent",
            bold=True,
        )
        add_text(
            slide,
            label,
            x + 0.30,
            3.20,
            3.1,
            0.45,
            size=20,
            color="primary",
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.30,
            3.90,
            3.1,
            1.2,
            size=14,
            color="secondary",
        )
    add_footer(slide, 4)

    # 5 — Layout 4: Minimal quote
    slide = prs.slides.add_slide(blank)
    set_background(slide, "bg")
    add_box(slide, 0.833, 2.30, 0.04, 2.0, fill="accent", line="accent", radius=False)
    add_text(
        slide,
        "“We don’t need more tools.\nWe need fewer repeated answers.”",
        1.25,
        2.20,
        10.8,
        2.2,
        size=36,
        color="primary",
        bold=True,
    )
    add_text(
        slide,
        "Support lead, fictional retail client — intake interview",
        1.25,
        4.85,
        8.0,
        0.3,
        size=14,
        color="secondary",
    )
    add_footer(slide, 5)

    # 6 — Closing action slide
    slide = prs.slides.add_slide(blank)
    set_background(slide, "bg")
    add_text(
        slide,
        "Decision: approve an 8-week pilot",
        0.833,
        0.72,
        10.0,
        0.7,
        size=33,
        color="primary",
        bold=True,
    )
    actions = [
        ("01", "Confirm scope", "Pick one workflow and name its owner."),
        ("02", "Open discovery", "Share approved examples and access limits."),
    ]
    for i, (number, label, body) in enumerate(actions):
        x = 0.833 + i * 5.95
        add_box(slide, x, 2.00, 5.67, 1.70, fill="card", no_line=True)
        add_text(
            slide,
            number,
            x + 0.32,
            2.28,
            0.8,
            0.4,
            size=18,
            color="accent",
            bold=True,
        )
        add_text(
            slide,
            label,
            x + 1.15,
            2.28,
            3.4,
            0.4,
            size=18,
            color="primary",
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 1.15,
            2.85,
            4.2,
            0.6,
            size=14,
            color="secondary",
        )
    add_box(slide, 0.833, 5.55, 11.67, 0.62, fill="accent", line="accent", corner_ratio=0.12)
    add_text(
        slide,
        "The ask: approve discovery. No production change yet.",
        1.15,
        5.55,
        11.0,
        0.62,
        size=15,
        color="white",
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, 6)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> int:
    parser = argparse.ArgumentParser()
    default = Path(__file__).resolve().parent.parent / "tests" / "smoke-test.pptx"
    parser.add_argument("--output", type=Path, default=default)
    args = parser.parse_args()
    make_deck(args.output)
    print(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
