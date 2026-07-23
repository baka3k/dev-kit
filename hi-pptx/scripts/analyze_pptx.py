#!/usr/bin/env python3
"""Inspect PPTX structure and styling without reproducing slide text."""

from __future__ import annotations

import argparse
import json
import re
import zipfile
from collections import Counter
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation


EMU_PER_INCH = 914400
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def theme_summary(path: Path) -> dict:
    with zipfile.ZipFile(path) as archive:
        try:
            root = ET.fromstring(archive.read("ppt/theme/theme1.xml"))
        except KeyError:
            return {"major_font": None, "minor_font": None, "colors": {}}

    ns = {"a": DRAWING_NS}
    major = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
    minor = root.find(".//a:fontScheme/a:minorFont/a:latin", ns)
    colors = {}
    scheme = root.find(".//a:clrScheme", ns)
    if scheme is not None:
        for slot in scheme:
            child = next(iter(slot), None)
            if child is not None:
                colors[slot.tag.rsplit("}", 1)[-1]] = (
                    child.attrib.get("val") or child.attrib.get("lastClr")
                )
    return {
        "major_font": major.attrib.get("typeface") if major is not None else None,
        "minor_font": minor.attrib.get("typeface") if minor is not None else None,
        "colors": colors,
    }


def font_color(run):
    try:
        rgb = run.font.color.rgb
        return str(rgb) if rgb else None
    except (AttributeError, TypeError):
        return None


def analyze(path: Path) -> dict:
    prs = Presentation(str(path))
    fonts = Counter()
    sizes = Counter()
    colors = Counter()
    layouts = Counter()
    slide_rows = []

    for index, slide in enumerate(prs.slides, start=1):
        layouts[slide.slide_layout.name] += 1
        words = text_shapes = pictures = tables = charts = 0
        for shape in iter_shapes(slide.shapes):
            if str(shape.shape_type).startswith("PICTURE"):
                pictures += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                text_shapes += 1
                words += len(re.findall(r"\b[\w'-]+\b", text))
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if run.font.size:
                        sizes[round(run.font.size.pt, 1)] += 1
                    color = font_color(run)
                    if color:
                        colors[color] += 1
        slide_rows.append(
            {
                "slide": index,
                "layout": slide.slide_layout.name,
                "words": words,
                "text_shapes": text_shapes,
                "pictures": pictures,
                "tables": tables,
                "charts": charts,
                "shape_count": len(slide.shapes),
            }
        )

    return {
        "file": path.name,
        "slide_count": len(prs.slides),
        "slide_size_inches": [
            round(prs.slide_width / EMU_PER_INCH, 3),
            round(prs.slide_height / EMU_PER_INCH, 3),
        ],
        "master_count": len(prs.slide_masters),
        "layouts": dict(layouts),
        "theme": theme_summary(path),
        "explicit_fonts": fonts.most_common(15),
        "explicit_font_sizes_pt": sizes.most_common(20),
        "explicit_text_colors": colors.most_common(15),
        "slides": slide_rows,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    if not args.pptx.is_file():
        parser.error(f"File not found: {args.pptx}")
    result = analyze(args.pptx)
    payload = json.dumps(result, ensure_ascii=False, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(payload + "\n", encoding="utf-8")
    else:
        print(payload)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
