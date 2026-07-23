#!/usr/bin/env python3
"""Run conservative structural and text-fit heuristics on a PPTX."""

from __future__ import annotations

import argparse
import json
import math
import re
from collections import Counter
from pathlib import Path

from pptx import Presentation


EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def font_sizes(shape) -> list[float]:
    values = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                values.append(run.font.size.pt)
    return values


def explicit_fonts(shape) -> list[str]:
    values = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                values.append(run.font.name)
    return values


def estimated_text_height(shape) -> float:
    frame = shape.text_frame
    width_pt = max(
        1.0,
        (shape.width - frame.margin_left - frame.margin_right) / EMU_PER_PT,
    )
    total = 0.0
    for paragraph in frame.paragraphs:
        sizes = [
            run.font.size.pt
            for run in paragraph.runs
            if run.font.size is not None
        ]
        size = max(sizes) if sizes else 16.0
        text = paragraph.text or " "
        chars_per_line = max(8, int(width_pt / max(4.0, size * 0.44)))
        lines = sum(
            max(1, math.ceil(len(segment) / chars_per_line))
            for segment in text.splitlines() or [text]
        )
        total += lines * size * 1.08
    return total


def overlap_ratio(a, b) -> float:
    left = max(a.left, b.left)
    top = max(a.top, b.top)
    right = min(a.left + a.width, b.left + b.width)
    bottom = min(a.top + a.height, b.top + b.height)
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    smaller = min(a.width * a.height, b.width * b.height)
    return float(intersection / smaller) if smaller else 0.0


def issue(level: str, slide: int, code: str, message: str) -> dict:
    return {"level": level, "slide": slide, "code": code, "message": message}


def lint(path: Path) -> dict:
    prs = Presentation(str(path))
    findings = []
    all_fonts = Counter()

    for slide_index, slide in enumerate(prs.slides, start=1):
        text_shapes = []
        word_count = 0
        for shape_index, shape in enumerate(iter_shapes(slide.shapes), start=1):
            right = shape.left + shape.width
            bottom = shape.top + shape.height
            tolerance = int(0.02 * EMU_PER_INCH)
            if (
                shape.left < -tolerance
                or shape.top < -tolerance
                or right > prs.slide_width + tolerance
                or bottom > prs.slide_height + tolerance
            ):
                findings.append(
                    issue(
                        "error",
                        slide_index,
                        "out-of-bounds",
                        f"Shape {shape_index} extends beyond the slide.",
                    )
                )

            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            text_shapes.append((shape_index, shape))
            word_count += len(re.findall(r"\b[\w'-]+\b", text))
            sizes = font_sizes(shape)
            fonts = explicit_fonts(shape)
            all_fonts.update(fonts)
            if sizes and min(sizes) < 9:
                findings.append(
                    issue(
                        "warning",
                        slide_index,
                        "tiny-text",
                        f"Shape {shape_index} contains text below 9 pt.",
                    )
                )

            available_pt = max(
                1.0,
                (
                    shape.height
                    - shape.text_frame.margin_top
                    - shape.text_frame.margin_bottom
                )
                / EMU_PER_PT,
            )
            estimate = estimated_text_height(shape)
            if estimate > available_pt * 1.22:
                findings.append(
                    issue(
                        "error",
                        slide_index,
                        "probable-overflow",
                        (
                            f"Shape {shape_index} estimated text height "
                            f"{estimate:.1f} pt exceeds {available_pt:.1f} pt."
                        ),
                    )
                )
            elif estimate > available_pt * 1.02:
                findings.append(
                    issue(
                        "warning",
                        slide_index,
                        "tight-text",
                        f"Shape {shape_index} is close to its estimated text limit.",
                    )
                )

        if word_count > 120:
            findings.append(
                issue(
                    "warning",
                    slide_index,
                    "dense-slide",
                    f"Slide contains approximately {word_count} words.",
                )
            )
        if len(text_shapes) > 18:
            findings.append(
                issue(
                    "warning",
                    slide_index,
                    "fragmented-slide",
                    f"Slide contains {len(text_shapes)} separate text shapes.",
                )
            )

        for left_pos, (left_index, left_shape) in enumerate(text_shapes):
            for right_index, right_shape in text_shapes[left_pos + 1 :]:
                if overlap_ratio(left_shape, right_shape) > 0.18:
                    findings.append(
                        issue(
                            "warning",
                            slide_index,
                            "text-overlap",
                            (
                                f"Text shapes {left_index} and {right_index} "
                                "overlap substantially."
                            ),
                        )
                    )

    if len(all_fonts) > 3:
        findings.append(
            issue(
                "warning",
                0,
                "font-variety",
                f"Deck explicitly uses {len(all_fonts)} fonts: {', '.join(all_fonts)}.",
            )
        )

    errors = sum(item["level"] == "error" for item in findings)
    warnings = sum(item["level"] == "warning" for item in findings)
    return {
        "file": path.name,
        "slide_count": len(prs.slides),
        "errors": errors,
        "warnings": warnings,
        "findings": findings,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()

    if not args.pptx.is_file():
        parser.error(f"File not found: {args.pptx}")
    result = lint(args.pptx)
    payload = json.dumps(result, ensure_ascii=False, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(payload + "\n", encoding="utf-8")
    print(f"{result['errors']} errors, {result['warnings']} warnings")
    if not args.output:
        print(payload)
    return 2 if result["errors"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
